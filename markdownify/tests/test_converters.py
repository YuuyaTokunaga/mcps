from __future__ import annotations

import base64
from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from markdownify_app.converters import csv_converter, docx_converter, excel_converter, pdf_converter, pptx_converter


def test_convert_csv_truncates_rows_and_columns(tmp_path: Path) -> None:
    csv_path = tmp_path / "data.csv"
    rows = ["c1,c2,c3"] + ["1,2,3" for _ in range(5)]
    csv_path.write_text("\n".join(rows), encoding="utf-8")

    markdown, images, meta, warnings = csv_converter.convert_csv(csv_path, max_rows=2, max_cols=2)

    assert "rows truncated due to max_rows" in warnings
    assert "columns truncated due to max_cols" in warnings
    assert "c1" in markdown
    assert images == []
    assert meta["rows"] == 2
    assert meta["cols"] == 2


def test_convert_excel_limits_sheets_and_reports_merges(tmp_path: Path) -> None:
    xlsx_path = tmp_path / "book.xlsx"
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Data"
    ws1.append(["A", "B"])
    ws1.append(["1", "2"])
    ws1.merge_cells("A1:B1")
    ws2 = wb.create_sheet("Extra")
    ws2.append(["X", "Y"])
    wb.save(xlsx_path)

    markdown, images, meta, warnings = excel_converter.convert_excel(
        xlsx_path,
        tmp_path / "out",
        include_images=False,
        max_rows=10,
        max_cols=5,
        max_sheets=1,
    )

    assert "sheets truncated due to max_sheets" in warnings
    assert any("merged cells" in w for w in warnings)
    assert markdown.startswith("## Sheet: Data")
    assert images == []
    assert meta["sheets"] == 1


def test_convert_pdf_truncates_pages(tmp_path: Path) -> None:
    pdf_path = tmp_path / "doc.pdf"
    doc = fitz.open()
    for idx in range(3):
        page = doc.new_page()
        page.insert_text((72, 72), f"Hello {idx + 1}")
    doc.save(pdf_path)
    doc.close()

    markdown, images, meta, warnings = pdf_converter.convert_pdf(
        pdf_path,
        tmp_path / "out",
        include_images=False,
        max_pages=2,
        render_dpi=50,
    )

    assert "pages truncated due to max_pages" in warnings
    assert markdown.startswith("## Page 1")
    assert images == []
    assert meta["pages_processed"] == 2


def test_convert_docx_extracts_headings_paragraphs_and_tables(tmp_path: Path) -> None:
    docx_path = tmp_path / "doc.docx"
    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph("Hello world")

    png_bytes = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMB/6X8bGQAAAAASUVORK5CYII="
    )
    png_path = tmp_path / "one.png"
    png_path.write_bytes(png_bytes)
    doc.add_picture(str(png_path))

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    doc.save(docx_path)

    markdown, images, meta, warnings = docx_converter.convert_docx(docx_path, max_paragraphs=100)

    assert "# Title" in markdown
    assert "Hello world" in markdown
    assert "| A | B |" in markdown
    assert images == []
    assert meta["paragraphs_processed"] >= 2
    assert meta["tables_processed"] == 1
    assert meta["figures"]["has_any"] is True
    assert meta["figures"]["inline_images"] >= 1
    assert warnings == []


def test_convert_pptx_extracts_slide_text(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])

    # 図形
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(1), Inches(1))

    # 画像
    png_bytes = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMB/6X8bGQAAAAASUVORK5CYII="
    )
    png_path = tmp_path / "one.png"
    png_path.write_bytes(png_bytes)
    slide.shapes.add_picture(str(png_path), Inches(2), Inches(0.5), width=Inches(1), height=Inches(1))

    # チャート（最小）
    chart_data = ChartData()
    chart_data.categories = ["A", "B"]
    chart_data.add_series("S", (1, 2))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2), Inches(3), Inches(2), chart_data)

    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = tx_box.text_frame
    tf.text = "Hello"
    p = tf.add_paragraph()
    p.text = "World"
    pres.save(pptx_path)

    markdown, images, meta, warnings = pptx_converter.convert_pptx(pptx_path, max_slides=10)

    assert markdown.startswith("## Slide 1")
    assert "- Hello" in markdown
    assert "- World" in markdown
    assert images == []
    assert meta["slides_processed"] == 1
    assert meta["figures"]["has_any"] is True
    assert meta["figures"]["pictures"] >= 1
    assert meta["figures"]["charts"] >= 1
    assert meta["figures"]["shapes"] >= 1
    assert meta["figures"]["slides_with_figures"] == [1]
    assert warnings == []
