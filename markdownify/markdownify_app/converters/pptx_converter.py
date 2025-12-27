from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _table_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""

    col_count = max((len(r) for r in rows), default=0)
    normalized = [(r + [""] * (col_count - len(r)))[:col_count] for r in rows]

    header = normalized[0]
    body = normalized[1:] if len(normalized) > 1 else []

    parts: list[str] = []
    parts.append("| " + " | ".join(cell or "" for cell in header) + " |")
    parts.append("|" + " --- |" * col_count)
    for row in body:
        parts.append("| " + " | ".join(cell or "" for cell in row) + " |")
    return "\n".join(parts)


def convert_pptx(
    input_path: Path,
    *,
    max_slides: int = 200,
    max_shapes_per_slide: int = 500,
    max_text_lines: int = 20_000,
    max_table_rows: int = 5_000,
    max_table_cols: int = 100,
) -> tuple[str, list[Path], dict[str, object], list[str]]:
    warnings: list[str] = []
    markdown_parts: list[str] = []
    images: list[Path] = []

    pres = Presentation(input_path)

    slides = list(pres.slides)
    if len(slides) > max_slides:
        warnings.append("slides truncated due to max_slides")
        slides = slides[:max_slides]

    text_lines_emitted = 0

    picture_count = 0
    chart_count = 0
    smartart_count = 0
    shape_count = 0
    slides_with_figures: list[int] = []

    smartart_type = getattr(MSO_SHAPE_TYPE, "SMART_ART", None)
    diagram_type = getattr(MSO_SHAPE_TYPE, "DIAGRAM", None)

    figure_shape_types = [
        getattr(MSO_SHAPE_TYPE, "AUTO_SHAPE", None),
        getattr(MSO_SHAPE_TYPE, "FREEFORM", None),
        getattr(MSO_SHAPE_TYPE, "CANVAS", None),
        getattr(MSO_SHAPE_TYPE, "GROUP", None),
        getattr(MSO_SHAPE_TYPE, "LINE", None),
        getattr(MSO_SHAPE_TYPE, "CONNECTOR", None),
    ]
    figure_shape_types_set = {t for t in figure_shape_types if t is not None}

    for s_idx, slide in enumerate(slides, start=1):
        title = None
        try:
            if slide.shapes.title is not None:
                title = (slide.shapes.title.text or "").strip() or None
        except Exception:  # noqa: BLE001
            title = None

        heading = f"## Slide {s_idx}" + (f": {title}" if title else "")
        markdown_parts.append(heading)

        slide_has_figure = False

        for sh_idx, shape in enumerate(slide.shapes, start=1):
            if sh_idx > max_shapes_per_slide:
                warnings.append("shapes truncated due to max_shapes_per_slide")
                break

            if getattr(shape, "has_table", False):
                table = shape.table
                rows: list[list[str]] = []
                for r_idx, row in enumerate(table.rows):
                    if r_idx >= max_table_rows:
                        warnings.append("table rows truncated due to max_table_rows")
                        break
                    cells: list[str] = []
                    for c_idx, cell in enumerate(row.cells):
                        if c_idx >= max_table_cols:
                            warnings.append("table cols truncated due to max_table_cols")
                            break
                        cells.append((cell.text or "").strip().replace("\n", " "))
                    rows.append(cells)
                if rows:
                    markdown_parts.append(_table_to_markdown(rows))
                    markdown_parts.append("")
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                slide_has_figure = True
                continue

            if getattr(shape, "has_chart", False):
                chart_count += 1
                slide_has_figure = True
                continue

            if smartart_type is not None and shape.shape_type == smartart_type:
                smartart_count += 1
                slide_has_figure = True
                continue
            if diagram_type is not None and shape.shape_type == diagram_type:
                smartart_count += 1
                slide_has_figure = True
                continue

            if figure_shape_types_set and shape.shape_type in figure_shape_types_set:
                shape_count += 1
                slide_has_figure = True

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # グループ内テキストは深追いしない（最低限の安全策）
                continue

            if getattr(shape, "has_text_frame", False):
                frame = shape.text_frame
                for para in frame.paragraphs:
                    text = (para.text or "").strip()
                    if not text:
                        continue
                    prefix = "  " * int(getattr(para, "level", 0))
                    markdown_parts.append(f"{prefix}- {text}")
                    text_lines_emitted += 1
                    if text_lines_emitted >= max_text_lines:
                        warnings.append("text truncated due to max_text_lines")
                        break
                if text_lines_emitted >= max_text_lines:
                    break

        markdown_parts.append("")
        if slide_has_figure:
            slides_with_figures.append(s_idx)
        if text_lines_emitted >= max_text_lines:
            break

    markdown = "\n".join(markdown_parts).strip()
    meta: dict[str, object] = {
        "slides_processed": len(slides),
        "max_slides": max_slides,
        "max_shapes_per_slide": max_shapes_per_slide,
        "max_text_lines": max_text_lines,
        "max_table_rows": max_table_rows,
        "max_table_cols": max_table_cols,
        "figures": {
            "has_any": bool(picture_count or chart_count or smartart_count or shape_count),
            "pictures": picture_count,
            "charts": chart_count,
            "smartart": smartart_count,
            "shapes": shape_count,
            "slides_with_figures": slides_with_figures,
        },
    }
    return markdown, images, meta, warnings
